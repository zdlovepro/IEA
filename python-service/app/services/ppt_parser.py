"""
PPT 解析服务
============
使用 python-pptx 读取 .pptx 文件，逐页提取文本、备注、表格内容。
图片与公式阶段二暂以占位符处理，阶段五集成 OCR / LaTeX 识别后补充。

保存位置：python-service/app/services/ppt_parser.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches  # noqa: F401  # 保留引入，后续可能用于尺寸判断
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.schemas.parse import PageContent, ParseResult
from app.utils.logger import logger


# ============================================================
# 公式检测辅助常量
# ============================================================
# PPT 中公式通常以 OLE 嵌入对象形式存在，其 prog_id 包含以下关键字之一。
# 不同版本 Office 的 prog_id 略有差异，这里覆盖常见情况。
_FORMULA_PROG_ID_KEYWORDS: list[str] = [
    "Equation",           # Microsoft Equation 3.0 → prog_id 含 "Equation.3" 或 "Equation.DSMT"
    "MathType",           # MathType 插件 → prog_id 含 "MathType"
    "OMML",               # Office MathML（较少以 OLE 出现，但保留以防万一）
]


def _is_formula_ole(shape) -> bool:
    """
    判断一个形状是否为公式 OLE 嵌入对象。

    Parameters
    ----------
    shape : pptx.shapes.base.BaseShape
        幻灯片中的单个形状。

    Returns
    -------
    bool
        如果该形状是公式 OLE 对象则返回 True，否则返回 False。

    Notes
    -----
    python-pptx 对 OLE 对象的支持有限，通过检查 shape 底层 XML 元素
    中的 progId 属性来判断是否为公式。这种方式兼容 Equation 3.0、
    MathType 等主流公式插件。
    """
    # OLE 嵌入对象在 XML 中以 <p:oleObj> 元素存在，
    # 其 progId 属性标识了嵌入对象的类型。
    # 我们遍历 shape 的 XML 子树寻找该元素。
    try:
        # shape._element 是 lxml Element，可以用 iter 遍历所有子孙节点
        for elem in shape._element.iter():
            # 标签名可能带命名空间前缀，例如
            # {http://schemas.openxmlformats.org/presentationml/2006/main}oleObj
            tag_local = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            if tag_local == "oleObj":
                prog_id = elem.get("progId", "")
                for keyword in _FORMULA_PROG_ID_KEYWORDS:
                    if keyword.lower() in prog_id.lower():
                        return True
    except Exception:
        # 如果 XML 解析出错，保守地返回 False，不影响整体解析流程
        pass

    return False


def _is_math_shape(shape) -> bool:
    """
    判断一个形状是否包含 Office MathML（OMML）公式。

    Parameters
    ----------
    shape : pptx.shapes.base.BaseShape
        幻灯片中的单个形状。

    Returns
    -------
    bool
        如果该形状包含 OMML 数学公式元素则返回 True。

    Notes
    -----
    Office 2010+ 支持直接在文本框中嵌入 OMML 公式（而非 OLE 对象）。
    这些公式以 <a:math> 或 <m:oMath> 元素的形式存在于形状的 XML 中。
    此函数通过检测这些元素来识别内嵌公式。
    """
    try:
        for elem in shape._element.iter():
            tag_local = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            # <a14:m> 或 <m:oMath> 或 <m:oMathPara> 都表示数学公式
            if tag_local in ("oMath", "oMathPara", "m"):
                return True
    except Exception:
        pass

    return False


def parse_pptx(file_path: str, courseware_id: str) -> ParseResult:
    """
    解析 .pptx 文件，返回结构化的 ParseResult。

    处理逻辑
    --------
    1. 逐页遍历所有幻灯片
    2. 对每个 shape 按类型分类处理：
       - 文本框 / 自选图形 → 提取文本
       - 表格 → 逐行拼接单元格文本
       - 图片（shape_type == PICTURE）→ 记录图片占位标识
       - OLE 嵌入对象（Equation / MathType）→ 记录公式占位标识
       - 内嵌 OMML 公式 → 记录公式占位标识
    3. 提取演讲者备注

    Parameters
    ----------
    file_path : str
        .pptx 文件的绝对路径。
    courseware_id : str
        课件唯一标识（由 Java 后端传入）。

    Returns
    -------
    ParseResult
        包含每页文本、备注、图片占位符及公式占位符的解析结果。

    Raises
    ------
    FileNotFoundError
        文件路径不存在时抛出。
    ValueError
        文件非 .pptx 格式时抛出。
    """
    path = Path(file_path)

    # ---- 前置校验 ----
    if not path.exists():
        raise FileNotFoundError(f"文件不存在：{file_path}")
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"不支持的文件格式：{path.suffix}，期望 .pptx")

    logger.info("开始解析 PPT 文件：%s（courseware_id=%s）", file_path, courseware_id)

    prs = Presentation(str(path))
    pages: list[PageContent] = []

    for idx, slide in enumerate(prs.slides, start=1):
        # ---- 每页的收集容器 ----
        texts: list[str] = []
        image_placeholders: list[str] = []
        formula_placeholders: list[str] = []

        for shape in slide.shapes:
            # --------------------------------------------------
            # 1) 公式检测（优先于文本提取，因为公式对象可能同时
            #    含有 text_frame，但其文本无意义）
            # --------------------------------------------------

            # 1a) OLE 嵌入公式（Equation 3.0 / MathType）
            if _is_formula_ole(shape):
                formula_placeholders.append(
                    f"[公式：slide_{idx}_formula_{len(formula_placeholders) + 1}]"
                )
                logger.debug(
                    "检测到 OLE 公式对象：slide_%d, shape_id=%s",
                    idx, shape.shape_id,
                )
                # OLE 公式对象的文本无意义，跳过后续文本提取
                continue

            # 1b) 内嵌 OMML 公式（Office 2010+ 原生公式）
            if _is_math_shape(shape):
                formula_placeholders.append(
                    f"[公式：slide_{idx}_formula_{len(formula_placeholders) + 1}]"
                )
                logger.debug(
                    "检测到 OMML 公式：slide_%d, shape_id=%s",
                    idx, shape.shape_id,
                )
                # OMML 公式的文本通常是公式源码，对用户无意义，跳过
                continue

            # --------------------------------------------------
            # 2) 图片检测
            # --------------------------------------------------
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_placeholders.append(
                    f"[图片：slide_{idx}_img_{len(image_placeholders) + 1}]"
                )
                continue

            # --------------------------------------------------
            # 3) 文本框 / 自选图形中的文本
            # --------------------------------------------------
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        texts.append(para_text)

            # --------------------------------------------------
            # 4) 表格：逐行拼接单元格文本
            # --------------------------------------------------
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]
                    if row_texts:
                        texts.append(" | ".join(row_texts))

        # ---- 提取演讲者备注 ----
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        # ---- 组装单页结果 ----
        page = PageContent(
            page_index=idx,
            text="\n".join(texts),
            notes=notes_text,
            image_placeholders=image_placeholders,
            formula_placeholders=formula_placeholders,
        )
        pages.append(page)

    result = ParseResult(
        courseware_id=courseware_id,
        file_type="pptx",
        total_pages=len(pages),
        pages=pages,
    )

    logger.info(
        "PPT 解析完成：共 %d 页（courseware_id=%s）",
        result.total_pages,
        courseware_id,
    )
    return result
